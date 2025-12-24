#!/usr/bin/env python3
"""
将导出的幻灯片图片合并为 PowerPoint 文件

使用方法:
    conda run -n test python create_pptx.py

输出:
    - presentation.pptx (16:9 宽屏格式)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu

# 配置
IMAGES_DIR = Path(__file__).parent / "slides_images"
OUTPUT_FILE = Path(__file__).parent / "presentation.pptx"
TOTAL_SLIDES = 20

# 幻灯片尺寸 (16:9 宽屏)
SLIDE_WIDTH = Inches(13.333)  # 1920px at 144dpi
SLIDE_HEIGHT = Inches(7.5)    # 1080px at 144dpi


def create_presentation():
    """创建 PowerPoint 演示文稿"""
    
    # 创建演示文稿
    prs = Presentation()
    
    # 设置幻灯片尺寸为 16:9 宽屏
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # 获取空白布局
    blank_layout = prs.slide_layouts[6]  # 空白布局
    
    print(f"📁 图片目录: {IMAGES_DIR}")
    print(f"📄 输出文件: {OUTPUT_FILE}")
    print()
    
    # 添加每张幻灯片
    for i in range(1, TOTAL_SLIDES + 1):
        image_path = IMAGES_DIR / f"slide_{i:02d}.png"
        
        if not image_path.exists():
            print(f"⚠️  跳过: {image_path.name} 不存在")
            continue
        
        # 添加新幻灯片
        slide = prs.slides.add_slide(blank_layout)
        
        # 添加图片，填满整个幻灯片
        slide.shapes.add_picture(
            str(image_path),
            left=Emu(0),
            top=Emu(0),
            width=SLIDE_WIDTH,
            height=SLIDE_HEIGHT
        )
        
        print(f"✅ 添加幻灯片 {i:2d}/20: {image_path.name}")
    
    # 保存演示文稿
    prs.save(str(OUTPUT_FILE))
    print()
    print(f"🎉 完成! PPT 已保存至: {OUTPUT_FILE}")
    print(f"   文件大小: {OUTPUT_FILE.stat().st_size / 1024 / 1024:.2f} MB")


if __name__ == "__main__":
    print("=" * 60)
    print("📊 创建 PowerPoint 演示文稿")
    print("=" * 60)
    print()
    
    create_presentation()
